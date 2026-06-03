# -*- coding: utf-8 -*-
"""
Trình tạo slide thuyết trình (.pptx) cho đồ án:
"Hệ thống Giám sát và Dự đoán Chất lượng Không khí sử dụng Edge AI (Air Quality AIoT)"
Nhóm KANT - Trường Đại học Lạc Hồng.

Chạy:  py make_slides.py
Xuất:  SlideThuyetTrinh_AIoT.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from PIL import Image

BASE = os.path.dirname(os.path.abspath(__file__))
MEDIA = os.path.join(BASE, "Bia_mau_media")
OUT = os.path.join(BASE, "SlideThuyetTrinh_AIoT.pptx")

# ---------------- BẢNG MÀU (Dark / Glassmorphism) ----------------
BG_TOP    = RGBColor(0x0B, 0x16, 0x27)   # nền gradient trên
BG_BOT    = RGBColor(0x12, 0x2A, 0x45)   # nền gradient dưới
CARD      = RGBColor(0x16, 0x27, 0x3E)   # mặt thẻ
CARD2     = RGBColor(0x1C, 0x33, 0x52)   # mặt thẻ nhấn
LINE      = RGBColor(0x2B, 0x44, 0x66)   # viền

CYAN   = RGBColor(0x22, 0xD3, 0xEE)
TEAL   = RGBColor(0x2D, 0xD4, 0xBF)
GREEN  = RGBColor(0x34, 0xD3, 0x99)
AMBER  = RGBColor(0xFB, 0xBF, 0x24)
RED    = RGBColor(0xF8, 0x71, 0x71)
BLUE   = RGBColor(0x60, 0xA5, 0xFA)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)

WHITE  = RGBColor(0xF1, 0xF5, 0xF9)
MUTED  = RGBColor(0x96, 0xA8, 0xC0)
DARKTX = RGBColor(0x0B, 0x16, 0x27)

FONT   = "Segoe UI"
FONT_B = "Segoe UI Semibold"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------------- HÀM TRỢ GIÚP ----------------
def _no_line(shape):
    shape.line.fill.background()


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE,
             line=None, line_w=Pt(1), radius=None, shadow=False):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    _solid(sp, color)
    if line is None:
        _no_line(sp)
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def gradient_bg(slide):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    _no_line(sp)
    sp.shadow.inherit = False
    try:
        sp.fill.gradient()
        stops = sp.fill.gradient_stops
        stops[0].position = 0.0
        stops[0].color.rgb = BG_TOP
        stops[1].position = 1.0
        stops[1].color.rgb = BG_BOT
        try:
            sp.fill.gradient_angle = 60.0
        except Exception:
            pass
    except Exception:
        _solid(sp, BG_TOP)
    return sp


def set_text(tf, runs, align=PP_ALIGN.LEFT, anchor=None, wrap=True):
    """runs: list of paragraphs; mỗi paragraph là list các (text, size, color, bold, italic)."""
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    if anchor is not None:
        tf.vertical_anchor = anchor
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = para.get("align", align)
        if "space_after" in para:
            p.space_after = para["space_after"]
        if "space_before" in para:
            p.space_before = para["space_before"]
        if "line" in para:
            p.line_spacing = para["line"]
        for (text, size, color, bold, *rest) in para["runs"]:
            italic = rest[0] if rest else False
            fname = rest[1] if len(rest) > 1 else FONT
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = fname
    return tf


def textbox(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.margin_left = 0
    tb.text_frame.margin_right = 0
    tb.text_frame.margin_top = 0
    tb.text_frame.margin_bottom = 0
    set_text(tb.text_frame, runs, align=align, anchor=anchor, wrap=wrap)
    return tb


def base_slide(section_no, section_name, title, accent=CYAN, subtitle=None):
    slide = prs.slides.add_slide(BLANK)
    gradient_bg(slide)
    # vệt sáng góc phải trên (trang trí)
    glow = add_rect(slide, SW - Inches(3.2), Inches(-1.4), Inches(4.6), Inches(4.6),
                    CARD2, shape=MSO_SHAPE.OVAL)
    glow.fill.fore_color.rgb = CARD2
    # chip section
    chip_w = Inches(0.62 + 0.12 * len(section_name))
    chip = add_rect(slide, Inches(0.7), Inches(0.55), chip_w, Inches(0.42),
                    accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    textbox(slide, Inches(0.7), Inches(0.55), chip_w, Inches(0.42),
            [{"runs": [("%02d  %s" % (section_no, section_name.upper()), 11.5, DARKTX, True)],
              "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    # tiêu đề
    textbox(slide, Inches(0.7), Inches(1.08), Inches(11.9), Inches(0.95),
            [{"runs": [(title, 30, WHITE, True, False, FONT_B)]}],
            anchor=MSO_ANCHOR.MIDDLE)
    # gạch nhấn dưới tiêu đề
    add_rect(slide, Inches(0.72), Inches(1.98), Inches(1.15), Pt(3.2), accent,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    if subtitle:
        textbox(slide, Inches(0.72), Inches(2.06), Inches(11.8), Inches(0.45),
                [{"runs": [(subtitle, 13.5, MUTED, False, True)]}])
    # footer
    add_rect(slide, Inches(0.7), SH - Inches(0.52), Inches(11.95), Pt(0.75), LINE)
    textbox(slide, Inches(0.7), SH - Inches(0.48), Inches(9), Inches(0.32),
            [{"runs": [("Hệ thống Giám sát Chất lượng Không khí AIoT  •  Nhóm KANT", 9.5, MUTED, False)]}],
            anchor=MSO_ANCHOR.MIDDLE)
    pageno = len(prs.slides._sldIdLst)
    textbox(slide, Inches(11.6), SH - Inches(0.48), Inches(1.05), Inches(0.32),
            [{"runs": [("%02d" % pageno, 11, accent, True)]}],
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def bullets(slide, x, y, w, h, items, marker=CYAN, size=15, gap=Pt(9),
            text_color=WHITE, anchor=MSO_ANCHOR.TOP, line=1.05):
    """items: list of (marker_text, body) hoặc chuỗi. body có thể chứa **bold** ở đầu."""
    paras = []
    for it in items:
        if isinstance(it, tuple):
            mk, body = it
        else:
            mk, body = "●", it
        runs = [(mk + "  ", size, marker, True)]
        # tách phần in đậm dẫn đầu nếu có dạng "Tiêu đề: phần còn lại"
        if isinstance(body, list):
            for seg in body:
                runs.append(seg)
        else:
            runs.append((body, size, text_color, False))
        paras.append({"runs": runs, "space_after": gap, "line": line})
    textbox(slide, x, y, w, h, paras, anchor=anchor)


def card(slide, x, y, w, h, fill=CARD, line=LINE, radius=0.06):
    return add_rect(slide, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
                    line=line, line_w=Pt(1), radius=radius)


def number_badge(slide, x, y, d, num, color):
    b = add_rect(slide, x, y, d, d, color, shape=MSO_SHAPE.OVAL)
    textbox(slide, x, y, d, d, [{"runs": [(str(num), 16, DARKTX, True)], "align": PP_ALIGN.CENTER}],
            anchor=MSO_ANCHOR.MIDDLE)
    return b


def fit_image(slide, path, x, y, w, h):
    im = Image.open(path)
    iw, ih = im.size
    ar = iw / ih
    box = w / h
    if ar > box:
        nw = w
        nh = int(w / ar)
    else:
        nh = h
        nw = int(h * ar)
    nx = x + (w - nw) // 2
    ny = y + (h - nh) // 2
    return slide.shapes.add_picture(path, nx, ny, nw, nh)


# ============================================================
# SLIDE 1 — TRANG BÌA
# ============================================================
def slide_cover():
    slide = prs.slides.add_slide(BLANK)
    gradient_bg(slide)
    # khối trang trí
    add_rect(slide, SW - Inches(4.3), Inches(-1.8), Inches(6), Inches(6),
             CARD2, shape=MSO_SHAPE.OVAL)
    add_rect(slide, Inches(-1.6), SH - Inches(3.2), Inches(5), Inches(5),
             CARD, shape=MSO_SHAPE.OVAL)
    add_rect(slide, 0, 0, Inches(0.18), SH, CYAN)
    # logo trường
    logo = os.path.join(MEDIA, "media", "image2.jpg")
    if os.path.exists(logo):
        fit_image(slide, logo, Inches(0.75), Inches(0.55), Inches(2.6), Inches(1.4))
    textbox(slide, Inches(0.78), Inches(2.0), Inches(8), Inches(0.8),
            [{"runs": [("TRƯỜNG ĐẠI HỌC LẠC HỒNG", 14, CYAN, True)], "space_after": Pt(2)},
             {"runs": [("KHOA CÔNG NGHỆ THÔNG TIN", 12.5, MUTED, True)]}])
    # nhãn loại tài liệu
    chip = add_rect(slide, Inches(0.8), Inches(2.95), Inches(3.1), Inches(0.4),
                    GREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    textbox(slide, Inches(0.8), Inches(2.95), Inches(3.1), Inches(0.4),
            [{"runs": [("BÁO CÁO ĐỒ ÁN MÔN HỌC", 11, DARKTX, True)], "align": PP_ALIGN.CENTER}],
            anchor=MSO_ANCHOR.MIDDLE)
    # tiêu đề chính
    textbox(slide, Inches(0.78), Inches(3.5), Inches(11.6), Inches(2.0),
            [{"runs": [("HỆ THỐNG GIÁM SÁT & DỰ ĐOÁN", 35, WHITE, True, False, FONT_B)],
              "space_after": Pt(2), "line": 1.02},
             {"runs": [("CHẤT LƯỢNG KHÔNG KHÍ ", 35, WHITE, True, False, FONT_B),
                       ("SỬ DỤNG EDGE AI", 35, CYAN, True, False, FONT_B)],
              "line": 1.02, "space_before": Pt(2)}])
    textbox(slide, Inches(0.8), Inches(5.25), Inches(11), Inches(0.5),
            [{"runs": [("Air Quality AIoT  •  ESP32-S3  •  LSTM TensorFlow Lite Micro  •  Firebase", 14, TEAL, False, True)]}])
    # thông tin nhóm
    info = [
        ("GVHD", "Hồ Ngọc Hoàng Long"),
        ("Lớp", "23CT111"),
        ("Nhóm", "KANT  —  Đồng Nai, 2026"),
    ]
    yy = Inches(5.95)
    textbox(slide, Inches(0.8), yy, Inches(7), Inches(1.2),
            [{"runs": [("GVHD:  ", 13, CYAN, True), ("Hồ Ngọc Hoàng Long", 13, WHITE, False),
                       ("        Lớp:  ", 13, CYAN, True), ("23CT111", 13, WHITE, False)],
              "space_after": Pt(4)},
             {"runs": [("SVTH:  ", 13, CYAN, True),
                       ("Tống Anh Kha  •  Nguyễn Thanh Nam  •  Đinh Thị Quỳnh Anh  •  Nguyễn Ngô Tôn Thái",
                        12.5, WHITE, False)]}])
    return slide


# ============================================================
# SLIDE 2 — NỘI DUNG TRÌNH BÀY
# ============================================================
def slide_agenda():
    slide = base_slide(0, "Mục lục", "Nội dung trình bày", accent=CYAN)
    items = [
        ("Giới thiệu", "Đặt vấn đề, mục tiêu & phạm vi nghiên cứu", GREEN),
        ("Thiết kế hệ thống", "Kiến trúc 3 tầng, phần cứng & sơ đồ đấu nối", BLUE),
        ("Cơ sở lý thuyết", "Lọc nhiễu Z-Score+EMA, mạng LSTM, lượng tử hóa INT8", PURPLE),
        ("Triển khai Firmware", "Luồng hoạt động & tối ưu trên ESP32-S3", AMBER),
        ("Web Dashboard", "Giao diện Glassmorphism đồng bộ thời gian thực", TEAL),
        ("Kết quả & Kết luận", "Thực nghiệm, đánh giá và hướng phát triển", RED),
    ]
    cols = 2
    cw = Inches(5.85)
    ch = Inches(1.32)
    x0 = Inches(0.72)
    y0 = Inches(2.55)
    gx = Inches(0.55)
    gy = Inches(0.28)
    for i, (t, d, c) in enumerate(items):
        r = i // cols
        col = i % cols
        x = x0 + col * (cw + gx)
        y = y0 + r * (ch + gy)
        card(slide, x, y, cw, ch, fill=CARD, line=LINE, radius=0.09)
        add_rect(slide, x, y, Inches(0.12), ch, c, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        number_badge(slide, x + Inches(0.32), y + Inches(0.34), Inches(0.62), i + 1, c)
        textbox(slide, x + Inches(1.18), y + Inches(0.2), cw - Inches(1.35), ch - Inches(0.3),
                [{"runs": [(t, 16.5, WHITE, True, False, FONT_B)], "space_after": Pt(3)},
                 {"runs": [(d, 11.5, MUTED, False)], "line": 1.0}],
                anchor=MSO_ANCHOR.MIDDLE)
    return slide


# ============================================================
# SLIDE 3 — ĐẶT VẤN ĐỀ
# ============================================================
def slide_problem():
    slide = base_slide(1, "Giới thiệu", "Đặt vấn đề — Ô nhiễm bụi mịn PM2.5",
                       accent=GREEN)
    # thẻ trái: tác hại
    card(slide, Inches(0.72), Inches(2.45), Inches(6.0), Inches(4.35), fill=CARD)
    textbox(slide, Inches(1.02), Inches(2.65), Inches(5.5), Inches(0.5),
            [{"runs": [("🫁  Mối đe dọa sức khỏe", 16, RED, True)]}])
    bullets(slide, Inches(1.02), Inches(3.25), Inches(5.45), Inches(3.4), [
        ("●", [("PM2.5", 14, WHITE, True), (" — hạt đường kính ≤ 2,5 µm, xuyên qua mũi họng, vào sâu phế nang và mạch máu.", 14, WHITE, False)]),
        ("●", [("Gây bệnh hô hấp mãn tính, tim mạch, ung thư phổi (theo WHO).", 14, WHITE, False)]),
        ("●", [("Ngưỡng an toàn WHO 2021: ", 14, WHITE, False), ("≤ 5 µg/m³/năm", 14, AMBER, True),
               (" — đô thị Việt Nam thường vượt nhiều lần.", 14, WHITE, False)]),
    ], marker=RED, size=14)
    # thẻ phải: hạn chế hiện trạng
    card(slide, Inches(6.92), Inches(2.45), Inches(5.7), Inches(4.35), fill=CARD)
    textbox(slide, Inches(7.22), Inches(2.65), Inches(5.2), Inches(0.5),
            [{"runs": [("⚠  Hạn chế của giải pháp hiện nay", 16, AMBER, True)]}])
    bullets(slide, Inches(7.22), Inches(3.25), Inches(5.15), Inches(3.4), [
        ("●", [("Trạm quan trắc ", 13.5, WHITE, False), ("thụ động", 13.5, AMBER, True),
               (": chỉ đo tức thời, không dự báo, không cảnh báo sớm.", 13.5, WHITE, False)]),
        ("●", [("Mô hình ", 13.5, WHITE, False), ("Cloud-only", 13.5, AMBER, True),
               (": độ trễ cao, tốn băng thông, mất cảnh báo khi rớt mạng.", 13.5, WHITE, False)]),
        ("●", [("Cảm biến bụi quang giá rẻ rất ", 13.5, WHITE, False),
               ("nhạy với gai nhiễu", 13.5, AMBER, True),
               (" (khói, bụi quét nhà, dị vật) → dự báo sai lệch.", 13.5, WHITE, False)]),
    ], marker=AMBER, size=13.5)
    # dải kết luận
    band = card(slide, Inches(0.72), Inches(6.92), Inches(11.9), Inches(0.0), fill=CARD)
    return slide


# ============================================================
# SLIDE 4 — GIẢI PHÁP ĐỀ XUẤT
# ============================================================
def slide_solution():
    slide = base_slide(1, "Giới thiệu", "Giải pháp đề xuất — Edge AI tại biên",
                       accent=GREEN,
                       subtitle="Nhúng trực tiếp thuật toán lọc nhiễu và mô hình AI lên vi điều khiển, hoạt động cả khi mất Internet")
    feats = [
        ("Đo đa chỉ số", "5 chỉ số môi trường: PM2.5, khí/VOCs, nhiệt độ, độ ẩm, áp suất — hiển thị tại chỗ trên LCD IPS.", BLUE, "📟"),
        ("Lọc nhiễu thông minh", "Bộ lọc Z-Score động + EMA (α=0,2) triệt tiêu gai nhiễu đột biến trước khi vào mô hình AI.", PURPLE, "🧹"),
        ("Dự báo bằng Edge AI", "Mô hình LSTM lượng tử hóa INT8 chạy trên PSRAM, suy luận 216 ms, dự báo PM2.5 1 giờ tới.", CYAN, "🧠"),
        ("Giám sát từ xa", "Đồng bộ Firebase + Web Dashboard Glassmorphism, theo dõi mọi lúc trên PC & điện thoại.", TEAL, "☁"),
    ]
    cw = Inches(2.92)
    ch = Inches(3.9)
    x0 = Inches(0.72)
    y0 = Inches(2.75)
    gx = Inches(0.18)
    for i, (t, d, c, ic) in enumerate(feats):
        x = x0 + i * (cw + gx)
        card(slide, x, y0, cw, ch, fill=CARD)
        add_rect(slide, x, y0, cw, Inches(0.12), c, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        circle = add_rect(slide, x + Inches(0.35), y0 + Inches(0.4), Inches(0.95), Inches(0.95),
                          CARD2, shape=MSO_SHAPE.OVAL, line=c, line_w=Pt(1.5))
        textbox(slide, x + Inches(0.35), y0 + Inches(0.4), Inches(0.95), Inches(0.95),
                [{"runs": [(ic, 26, c, False)], "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, x + Inches(0.28), y0 + Inches(1.55), cw - Inches(0.56), Inches(0.7),
                [{"runs": [(t, 15.5, WHITE, True, False, FONT_B)]}], anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, x + Inches(0.28), y0 + Inches(2.2), cw - Inches(0.56), Inches(1.55),
                [{"runs": [(d, 12, MUTED, False)], "line": 1.08}])
    return slide


# ============================================================
# SLIDE 5 — MỤC TIÊU
# ============================================================
def slide_objectives():
    slide = base_slide(1, "Giới thiệu", "Mục tiêu của đồ án", accent=GREEN)
    goals = [
        ("Thiết kế phần cứng trạm đo đa chỉ số",
         "Đo đồng thời PM2.5 (Sharp GP2Y1014), khí/VOCs (MQ135), nhiệt độ – độ ẩm – áp suất (BME280), hiển thị trực tiếp trên LCD IPS ST7789.", BLUE),
        ("Lọc nhiễu Z-Score + EMA trên vi điều khiển",
         "Lập trình C++ bộ phát hiện dị thường ngay trên ESP32-S3, triệt tiêu gai nhiễu cục bộ, đảm bảo đầu vào sạch cho AI.", PURPLE),
        ("Nhúng mô hình LSTM INT8 chạy tại biên",
         "Huấn luyện LSTM trên Python, lượng tử hóa Post-Training INT8, nhúng lên PSRAM qua TensorFlow Lite Micro (< 250 ms/chu kỳ).", CYAN),
        ("Xây dựng hệ sinh thái Cloud & Web Dashboard",
         "Đồng bộ Firebase Realtime Database, giao diện Glassmorphism giám sát từ xa đa thiết bị qua HTTPS.", TEAL),
    ]
    x0 = Inches(0.72)
    y0 = Inches(2.5)
    cw = Inches(11.9)
    ch = Inches(1.0)
    gy = Inches(0.14)
    for i, (t, d, c) in enumerate(goals):
        y = y0 + i * (ch + gy)
        card(slide, x0, y, cw, ch, fill=CARD)
        number_badge(slide, x0 + Inches(0.28), y + Inches(0.2), Inches(0.6), i + 1, c)
        textbox(slide, x0 + Inches(1.15), y + Inches(0.13), cw - Inches(1.4), Inches(0.45),
                [{"runs": [(t, 15.5, WHITE, True, False, FONT_B)]}])
        textbox(slide, x0 + Inches(1.15), y + Inches(0.55), cw - Inches(1.4), Inches(0.42),
                [{"runs": [(d, 12, MUTED, False)], "line": 1.0}])
    return slide


# ============================================================
# SLIDE 6 — PHẠM VI NGHIÊN CỨU
# ============================================================
def slide_scope():
    slide = base_slide(1, "Giới thiệu", "Phạm vi & đối tượng nghiên cứu", accent=GREEN)
    # cột trái: phạm vi
    card(slide, Inches(0.72), Inches(2.5), Inches(6.0), Inches(4.3), fill=CARD)
    textbox(slide, Inches(1.02), Inches(2.7), Inches(5.5), Inches(0.5),
            [{"runs": [("✔  Trong phạm vi", 16, GREEN, True)]}])
    bullets(slide, Inches(1.02), Inches(3.3), Inches(5.45), Inches(3.3), [
        ("●", [("Thiết bị đơn lẻ, đặt ", 13.5, WHITE, False), ("indoor", 13.5, GREEN, True),
               (" (phòng ngủ, văn phòng) hoặc ", 13.5, WHITE, False), ("bán-outdoor", 13.5, GREEN, True),
               (" (ban công, hành lang).", 13.5, WHITE, False)]),
        ("●", [("Dữ liệu chuỗi thời gian thu thập ", 13.5, WHITE, False),
               ("5 giây/mẫu", 13.5, GREEN, True), (".", 13.5, WHITE, False)]),
        ("●", [("LSTM dùng cửa sổ trượt ", 13.5, WHITE, False), ("24 bước", 13.5, GREEN, True),
               (" để dự báo PM2.5 trong ", 13.5, WHITE, False), ("1 giờ tới", 13.5, GREEN, True),
               (".", 13.5, WHITE, False)]),
    ], marker=GREEN, size=13.5)
    # cột phải: giới hạn
    card(slide, Inches(6.92), Inches(2.5), Inches(5.7), Inches(4.3), fill=CARD)
    textbox(slide, Inches(7.22), Inches(2.7), Inches(5.2), Inches(0.5),
            [{"runs": [("✘  Ngoài phạm vi", 16, RED, True)]}])
    bullets(slide, Inches(7.22), Inches(3.3), Inches(5.15), Inches(3.3), [
        ("●", [("Chưa xây dựng mạng lưới nhiều trạm quan trắc phân tán.", 13.5, WHITE, False)]),
        ("●", [("Chưa phát triển ứng dụng di động native.", 13.5, WHITE, False)]),
        ("●", [("Chưa tích hợp pin sạc — dùng nguồn USB/Adapter cố định.", 13.5, WHITE, False)]),
    ], marker=RED, size=13.5)
    return slide


# ============================================================
# SLIDE 7 — KIẾN TRÚC TỔNG THỂ 3 TẦNG
# ============================================================
def slide_architecture():
    slide = base_slide(2, "Thiết kế hệ thống", "Kiến trúc tổng thể — Mô hình 3 tầng", accent=BLUE)
    tiers = [
        ("EDGE LAYER", "Tầng cảm biến & xử lý biên", CYAN,
         ["ESP32-S3 thu thập 3 cảm biến", "Lọc nhiễu Z-Score / EMA", "Suy luận AI (TFLite Micro)", "Điều khiển cảnh báo + LCD"]),
        ("CLOUD LAYER", "Tầng đám mây", AMBER,
         ["Firebase Realtime Database", "/sensor & /ai — mỗi 5 giây", "/history — mỗi 60 giây", "Server Timestamp của Google"]),
        ("PRESENTATION", "Tầng giao diện người dùng", TEAL,
         ["Web Dashboard HTML/CSS/JS", "Firebase Hosting + CDN", "Truy cập đa thiết bị (HTTPS)", "Glassmorphism + Chart.js"]),
    ]
    cw = Inches(3.72)
    ch = Inches(3.9)
    x0 = Inches(0.72)
    y0 = Inches(2.7)
    gx = Inches(0.45)
    for i, (tag, name, c, pts) in enumerate(tiers):
        x = x0 + i * (cw + gx)
        card(slide, x, y0, cw, ch, fill=CARD)
        head = add_rect(slide, x, y0, cw, Inches(0.92), c, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
        textbox(slide, x, y0 + Inches(0.1), cw, Inches(0.45),
                [{"runs": [(tag, 15, DARKTX, True)], "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, x, y0 + Inches(0.5), cw, Inches(0.36),
                [{"runs": [(name, 11, DARKTX, False, True)], "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
        bullets(slide, x + Inches(0.32), y0 + Inches(1.15), cw - Inches(0.6), Inches(2.6),
                [("▸", p) for p in pts], marker=c, size=12.5, gap=Pt(8))
        # mũi tên nối
        if i < 2:
            ar = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + cw + Inches(0.02),
                                        y0 + Inches(1.6), Inches(0.42), Inches(0.6))
            _solid(ar, c)
            _no_line(ar)
            ar.shadow.inherit = False
    return slide


# ============================================================
# SLIDE 8 — BOM
# ============================================================
def slide_bom():
    slide = base_slide(2, "Thiết kế hệ thống", "Danh mục phần cứng (Bill of Materials)", accent=BLUE)
    rows = [
        ("Linh kiện", "Thông số kỹ thuật", "Vai trò trong hệ thống"),
        ("ESP32-S3-WROOM-1 (N16R8)", "Dual-core LX7 240 MHz; 16 MB Flash; 8 MB PSRAM", "Nhân xử lý: lọc nhiễu, chạy AI, điều khiển, WiFi"),
        ("Sharp GP2Y1014AU0F", "Cảm biến bụi quang IR; ngõ ra Analog; 5 V", "Đo nồng độ PM2.5 sơ bộ (tương đối)"),
        ("MQ135 Gas Sensor", "Cảm biến bán dẫn; NH₃/NOx/CO₂/Benzen; ~150 mA", "Đo khí độc hại & VOCs, cảnh báo ô nhiễm khí"),
        ("ST7789 IPS LCD", "1,54\"; 240×240 px; giao tiếp SPI", "Hiển thị Dashboard tại chỗ"),
        ("Buzzer + LED đỏ", "Còi active 5 V (~2,3 kHz); LED 5 mm", "Cảnh báo âm thanh & ánh sáng"),
        ("MB102 Power Module", "5 V / 3,3 V; tối đa 700 mA; DC Barrel", "Cấp nguồn ổn định, chống sụt áp (brown-out)"),
    ]
    add_table(slide, rows, Inches(0.72), Inches(2.55), Inches(11.9), Inches(4.2),
              col_w=[Inches(3.1), Inches(4.6), Inches(4.2)], header_color=BLUE,
              font_size=11.5, header_size=12.5)
    return slide


def add_table(slide, rows, x, y, w, h, col_w, header_color, font_size=11.5,
              header_size=12.5, first_bold=True):
    nr = len(rows)
    nc = len(rows[0])
    gtable = slide.shapes.add_table(nr, nc, x, y, w, h).table
    gtable.first_row = False
    gtable.horz_banding = False
    for j, cw in enumerate(col_w):
        gtable.columns[j].width = cw
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = gtable.cell(i, j)
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_color
                col = DARKTX; bold = True; fs = header_size
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD if i % 2 == 1 else CARD2
                col = WHITE if j == 0 else MUTED
                bold = (j == 0 and first_bold)
                fs = font_size
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(fs); r.font.color.rgb = col
            r.font.bold = bold; r.font.name = FONT
            if i == 0 and j == 0:
                r.font.color.rgb = DARKTX
    return gtable


# ============================================================
# SLIDE 9 — PINOUT + RC FILTER
# ============================================================
def slide_pinout():
    slide = base_slide(2, "Thiết kế hệ thống", "Sơ đồ đấu nối GPIO & mạch bảo vệ", accent=BLUE)
    rows = [
        ("Ngoại vi", "Tín hiệu", "GPIO", "Giao thức"),
        ("BME280", "SDA / SCL", "8 / 9", "I2C Master"),
        ("Sharp GP2Y1014", "Vo / LED", "4 / 5", "ADC1 / Xung 10 ms"),
        ("MQ135", "AO", "1", "ADC1"),
        ("ST7789 LCD", "MOSI/SCLK/CS/DC/RST/BLK", "11/12/10/13/14/15", "SPI HW + PWM"),
        ("Buzzer", "Tín hiệu", "6", "PWM tone()"),
        ("LED cảnh báo", "Dương cực", "7", "Digital Output"),
        ("Nút BOOT", "Reset WiFi", "0", "Input Pull-up"),
    ]
    add_table(slide, rows, Inches(0.72), Inches(2.5), Inches(7.55), Inches(4.25),
              col_w=[Inches(1.95), Inches(2.55), Inches(1.65), Inches(1.4)],
              header_color=BLUE, font_size=11, header_size=11.5)
    # thẻ phải: mạch lọc RC
    card(slide, Inches(8.5), Inches(2.5), Inches(4.12), Inches(4.25), fill=CARD)
    textbox(slide, Inches(8.8), Inches(2.72), Inches(3.6), Inches(0.5),
            [{"runs": [("🔧  Mạch lọc RC thông thấp", 14.5, AMBER, True)]}])
    textbox(slide, Inches(8.8), Inches(3.25), Inches(3.55), Inches(0.9),
            [{"runs": [("Bảo vệ cảm biến bụi Sharp khỏi gai điện áp khi LED IR nháy xung.", 12, MUTED, False)], "line": 1.08}])
    # công thức
    eq = add_rect(slide, Inches(8.8), Inches(4.2), Inches(3.55), Inches(0.95),
                  CARD2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=AMBER, radius=0.1)
    textbox(slide, Inches(8.8), Inches(4.2), Inches(3.55), Inches(0.95),
            [{"runs": [("τ = R · C", 16, WHITE, True)], "align": PP_ALIGN.CENTER, "space_after": Pt(2)},
             {"runs": [("= 150 Ω × 220 µF = ", 12.5, MUTED, False), ("33 ms", 13.5, AMBER, True)],
              "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    bullets(slide, Inches(8.8), Inches(5.35), Inches(3.6), Inches(1.3), [
        ("▸", [("Lọc nhiễu cao tần > 30 Hz trên đường nguồn.", 11.5, WHITE, False)]),
        ("▸", [("Nguồn MB102 riêng cho MQ135 + Sharp, chung GND.", 11.5, WHITE, False)]),
    ], marker=AMBER, size=11.5, gap=Pt(6))
    return slide


# ============================================================
# SLIDE 10 — PIPELINE 3 TẦNG XỬ LÝ
# ============================================================
def slide_pipeline():
    slide = base_slide(3, "Cơ sở lý thuyết", "Pipeline xử lý dữ liệu — 3 tầng tích hợp", accent=PURPLE,
                       subtitle="Đóng góp kỹ thuật cốt lõi: dữ liệu thô được làm sạch, dự báo và đồng bộ liền mạch")
    stages = [
        ("Dữ liệu thô", "3 cảm biến\n5 giây/mẫu", LINE, "📡"),
        ("Lọc dị thường", "Z-Score + EMA\nchặn gai nhiễu", PURPLE, "🧹"),
        ("Chuẩn hóa + AI", "MinMaxScaler →\nLSTM INT8 (216 ms)", CYAN, "🧠"),
        ("Cảnh báo + Cloud", "Buzzer/LED +\nFirebase sync", GREEN, "🔔"),
    ]
    cw = Inches(2.72)
    ch = Inches(2.7)
    x0 = Inches(0.72)
    y0 = Inches(3.2)
    gx = Inches(0.5)
    for i, (t, d, c, ic) in enumerate(stages):
        x = x0 + i * (cw + gx)
        col = c if c != LINE else BLUE
        card(slide, x, y0, cw, ch, fill=CARD, line=col)
        textbox(slide, x, y0 + Inches(0.35), cw, Inches(0.7),
                [{"runs": [(ic, 30, col, False)], "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, x, y0 + Inches(1.15), cw, Inches(0.5),
                [{"runs": [(t, 15, WHITE, True, False, FONT_B)], "align": PP_ALIGN.CENTER}])
        textbox(slide, x, y0 + Inches(1.7), cw, Inches(0.9),
                [{"runs": [(d, 12, MUTED, False)], "align": PP_ALIGN.CENTER, "line": 1.05}])
        if i < 3:
            ar = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + cw + Inches(0.05),
                                        y0 + Inches(1.05), Inches(0.4), Inches(0.55))
            _solid(ar, col); _no_line(ar); ar.shadow.inherit = False
    # dải nhấn dưới
    card(slide, Inches(0.72), Inches(6.2), Inches(11.9), Inches(0.62), fill=CARD2, line=PURPLE)
    textbox(slide, Inches(0.72), Inches(6.2), Inches(11.9), Inches(0.62),
            [{"runs": [("Tất cả xử lý diễn ra ngay trên ESP32-S3 — thiết bị vẫn cảnh báo kể cả khi mất hoàn toàn kết nối Internet.",
                        13, WHITE, False, True)], "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    return slide


# ============================================================
# SLIDE 11 — CHUẨN HÓA DỮ LIỆU
# ============================================================
def slide_scaler():
    slide = base_slide(3, "Cơ sở lý thuyết", "Chuẩn hóa dữ liệu (MinMaxScaler)", accent=PURPLE)
    card(slide, Inches(0.72), Inches(2.5), Inches(6.0), Inches(4.3), fill=CARD)
    textbox(slide, Inches(1.02), Inches(2.7), Inches(5.5), Inches(0.9),
            [{"runs": [("LSTM rất nhạy với biên độ đầu vào. Mọi đặc trưng được đưa về dải [0, 1] thống nhất:", 13.5, WHITE, False)], "line": 1.1}])
    eq = add_rect(slide, Inches(1.02), Inches(3.65), Inches(5.4), Inches(0.95),
                  CARD2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=PURPLE, radius=0.1)
    textbox(slide, Inches(1.02), Inches(3.65), Inches(5.4), Inches(0.95),
            [{"runs": [("x_scaled = (x − x_min) / (x_max − x_min)", 15, WHITE, True)],
              "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    bullets(slide, Inches(1.02), Inches(4.85), Inches(5.5), Inches(1.9), [
        ("▸", [("Hằng số x_min/x_max ", 12.5, WHITE, False), ("đồng bộ tuyệt đối", 12.5, AMBER, True),
               (" giữa Python và firmware C++.", 12.5, WHITE, False)]),
        ("▸", [("Hàm scaleData() có ", 12.5, WHITE, False), ("clamping", 12.5, AMBER, True),
               (" xử lý giá trị NaN từ cảm biến lỗi.", 12.5, WHITE, False)]),
        ("▸", [("Inverse Scale phục hồi kết quả về µg/m³ sau suy luận.", 12.5, WHITE, False)]),
    ], marker=PURPLE, size=12.5, gap=Pt(8))
    # bảng dải đặc trưng
    rows = [
        ("Đặc trưng", "x_min", "x_max", "Đơn vị"),
        ("PM2.5", "0,0", "1000,0", "µg/m³"),
        ("Độ ẩm / Điểm sương", "−40,0", "30,0", "% RH"),
        ("Nhiệt độ", "−20,0", "45,0", "°C"),
        ("Áp suất", "990,0", "1050,0", "hPa"),
    ]
    add_table(slide, rows, Inches(6.95), Inches(2.85), Inches(5.65), Inches(3.0),
              col_w=[Inches(2.55), Inches(1.15), Inches(1.15), Inches(0.8)],
              header_color=PURPLE, font_size=12, header_size=12.5)
    textbox(slide, Inches(6.95), Inches(5.95), Inches(5.65), Inches(0.7),
            [{"runs": [("Khớp chính xác mảng FEATURE_MIN[] / FEATURE_MAX[] trong main.cpp.", 11.5, MUTED, False, True)], "line": 1.05}])
    return slide


# ============================================================
# SLIDE 12 — Z-SCORE + EMA
# ============================================================
def slide_zscore():
    slide = base_slide(3, "Cơ sở lý thuyết", "Lọc dị thường: Z-Score động + EMA", accent=PURPLE,
                       subtitle="Triệt tiêu gai nhiễu đột biến từ cảm biến quang trước khi đưa vào LSTM")
    # cột trái: công thức
    card(slide, Inches(0.72), Inches(2.75), Inches(5.85), Inches(4.05), fill=CARD)
    textbox(slide, Inches(1.0), Inches(2.92), Inches(5.3), Inches(0.4),
            [{"runs": [("Đường cơ sở thích nghi (EMA, α = 0,2)", 14, CYAN, True)]}])
    for i, (label, formula) in enumerate([
        ("EMA", "EMAₜ = α·xₜ + (1−α)·EMAₜ₋₁"),
        ("Phương sai", "Varₜ = α·(xₜ−EMAₜ₋₁)² + (1−α)·Varₜ₋₁"),
        ("Z-Score", "Zₜ = |xₜ − EMAₜ₋₁| / max(σₜ, 10)"),
    ]):
        yy = Inches(3.4) + i * Inches(0.92)
        box = add_rect(slide, Inches(1.0), yy, Inches(5.3), Inches(0.78),
                       CARD2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=LINE, radius=0.12)
        textbox(slide, Inches(1.15), yy + Inches(0.06), Inches(5.0), Inches(0.3),
                [{"runs": [(label, 10.5, PURPLE, True)]}])
        textbox(slide, Inches(1.15), yy + Inches(0.34), Inches(5.0), Inches(0.4),
                [{"runs": [(formula, 13, WHITE, True)]}])
    # cột phải: quy tắc
    card(slide, Inches(6.8), Inches(2.75), Inches(5.82), Inches(4.05), fill=CARD)
    textbox(slide, Inches(7.1), Inches(2.92), Inches(5.3), Inches(0.4),
            [{"runs": [("Quy tắc phán định 3σ", 14, AMBER, True)]}])
    # Z<=3
    g = add_rect(slide, Inches(7.1), Inches(3.4), Inches(5.2), Inches(1.45),
                 CARD2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=GREEN, radius=0.08)
    textbox(slide, Inches(7.3), Inches(3.52), Inches(4.9), Inches(1.25),
            [{"runs": [("Z ≤ 3,0  →  Bình thường ✔", 14, GREEN, True)], "space_after": Pt(3)},
             {"runs": [("Mẫu đi qua, cập nhật EMA & Var. Bao phủ 99,7% dữ liệu hợp lệ.", 12, WHITE, False)], "line": 1.05}])
    # Z>3
    r = add_rect(slide, Inches(7.1), Inches(5.0), Inches(5.2), Inches(1.6),
                 CARD2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=RED, radius=0.08)
    textbox(slide, Inches(7.3), Inches(5.12), Inches(4.9), Inches(1.4),
            [{"runs": [("Z > 3,0  →  Gai nhiễu ✘", 14, RED, True)], "space_after": Pt(3)},
             {"runs": [("Giá trị thô bị chặn, thay bằng EMAₜ₋₁. EMA & Var giữ nguyên để bảo toàn đường cơ sở (xác suất sai < 0,3%).", 12, WHITE, False)], "line": 1.05}])
    return slide


# ============================================================
# SLIDE 13 — LSTM LÝ THUYẾT
# ============================================================
def slide_lstm():
    slide = base_slide(3, "Cơ sở lý thuyết", "Mạng nơ-ron hồi quy LSTM", accent=PURPLE)
    card(slide, Inches(0.72), Inches(2.5), Inches(5.5), Inches(4.3), fill=CARD)
    textbox(slide, Inches(1.0), Inches(2.7), Inches(5.0), Inches(0.45),
            [{"runs": [("Vì sao chọn LSTM?", 15, CYAN, True)]}])
    bullets(slide, Inches(1.0), Inches(3.25), Inches(5.0), Inches(3.4), [
        ("▸", [("Dự báo PM2.5 là bài toán ", 13, WHITE, False), ("chuỗi thời gian", 13, CYAN, True),
               (" — giá trị tương lai phụ thuộc lịch sử quá khứ.", 13, WHITE, False)]),
        ("▸", [("Mạng truyền thẳng không có bộ nhớ; RNN cơ bản bị ", 13, WHITE, False),
               ("tiêu biến đạo hàm", 13, AMBER, True), (".", 13, WHITE, False)]),
        ("▸", [("LSTM (Hochreiter & Schmidhuber, 1997) dùng ", 13, WHITE, False),
               ("Cell State", 13, GREEN, True),
               (" — đường truyền gradient xuyên thời gian, học được phụ thuộc dài hạn.", 13, WHITE, False)]),
    ], marker=CYAN, size=13, gap=Pt(10))
    # cột phải: 3 cổng
    card(slide, Inches(6.4), Inches(2.5), Inches(6.2), Inches(4.3), fill=CARD)
    textbox(slide, Inches(6.7), Inches(2.7), Inches(5.6), Inches(0.45),
            [{"runs": [("Ba cổng kiểm soát thông tin", 15, AMBER, True)]}])
    gates = [
        ("Cổng Quên (fₜ)", "Quyết định thông tin dài hạn nào trong Cell State cũ cần xóa.", RED),
        ("Cổng Vào (iₜ)", "Xác định thông tin mới từ đầu vào sẽ được ghi vào ô nhớ.", GREEN),
        ("Cổng Ra (oₜ)", "Quyết định phần nào của ô nhớ xuất ra trạng thái ẩn hₜ.", BLUE),
    ]
    for i, (t, d, c) in enumerate(gates):
        yy = Inches(3.25) + i * Inches(1.12)
        box = add_rect(slide, Inches(6.7), yy, Inches(5.6), Inches(0.98),
                       CARD2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=c, radius=0.08)
        add_rect(slide, Inches(6.7), yy, Inches(0.1), Inches(0.98), c,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        textbox(slide, Inches(6.95), yy + Inches(0.1), Inches(5.2), Inches(0.4),
                [{"runs": [(t, 13.5, c, True)]}])
        textbox(slide, Inches(6.95), yy + Inches(0.48), Inches(5.2), Inches(0.45),
                [{"runs": [(d, 11.5, WHITE, False)], "line": 1.0}])
    return slide


# ============================================================
# SLIDE 14 — MÔ HÌNH LSTM NHÚNG
# ============================================================
def slide_model():
    slide = base_slide(3, "Cơ sở lý thuyết", "Cấu trúc mô hình LSTM nhúng", accent=PURPLE)
    layers = [
        ("INPUT", "[1, 24, 4]", "24 bước thời gian × 4 đặc trưng\n(PM2.5, Độ ẩm, Nhiệt độ, Áp suất)", BLUE),
        ("LSTM LAYER", "32 units", "4 cổng × 1152 trọng số = 4608\n+ 128 biases", PURPLE),
        ("DENSE OUTPUT", "1 neuron", "y_scaled ∈ [0,1]\n→ Inverse Scale → µg/m³", GREEN),
    ]
    cw = Inches(3.72)
    ch = Inches(2.85)
    x0 = Inches(0.72)
    y0 = Inches(2.75)
    gx = Inches(0.45)
    for i, (tag, val, d, c) in enumerate(layers):
        x = x0 + i * (cw + gx)
        card(slide, x, y0, cw, ch, fill=CARD, line=c)
        textbox(slide, x, y0 + Inches(0.25), cw, Inches(0.4),
                [{"runs": [(tag, 13, c, True)], "align": PP_ALIGN.CENTER}])
        textbox(slide, x, y0 + Inches(0.75), cw, Inches(0.6),
                [{"runs": [(val, 24, WHITE, True, False, FONT_B)], "align": PP_ALIGN.CENTER}])
        textbox(slide, x + Inches(0.25), y0 + Inches(1.65), cw - Inches(0.5), Inches(1.0),
                [{"runs": [(d, 12, MUTED, False)], "align": PP_ALIGN.CENTER, "line": 1.1}])
        if i < 2:
            ar = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + cw + Inches(0.02),
                                        y0 + Inches(1.15), Inches(0.42), Inches(0.55))
            _solid(ar, c); _no_line(ar); ar.shadow.inherit = False
    # dải dưới: circular buffer
    card(slide, Inches(0.72), Inches(5.95), Inches(11.9), Inches(0.85), fill=CARD2, line=PURPLE)
    textbox(slide, Inches(1.0), Inches(5.95), Inches(11.4), Inches(0.85),
            [{"runs": [("Bộ đệm vòng  ", 13, PURPLE, True),
                       ("int8_t input_buffer[24][4]", 13, CYAN, True),
                       ("  — mỗi chu kỳ 5 s dịch trái 1 bước, ghi mẫu mới nhất vào input_buffer[23] (mô phỏng cửa sổ trượt).",
                        13, WHITE, False)], "line": 1.05}],
            anchor=MSO_ANCHOR.MIDDLE)
    return slide


# ============================================================
# SLIDE 15 — LƯỢNG TỬ HÓA INT8
# ============================================================
def slide_quant():
    slide = base_slide(3, "Cơ sở lý thuyết", "Lượng tử hóa INT8 (Post-Training Quantization)", accent=PURPLE)
    card(slide, Inches(0.72), Inches(2.5), Inches(5.85), Inches(4.3), fill=CARD)
    textbox(slide, Inches(1.0), Inches(2.7), Inches(5.3), Inches(0.9),
            [{"runs": [("ESP32-S3 không có FPU vector, RAM nội bộ chỉ 320 KB → bắt buộc nén mô hình từ Float32 (4 byte) về INT8 (1 byte):", 13, WHITE, False)], "line": 1.1}])
    eq = add_rect(slide, Inches(1.0), Inches(3.75), Inches(5.3), Inches(0.85),
                  CARD2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=PURPLE, radius=0.1)
    textbox(slide, Inches(1.0), Inches(3.75), Inches(5.3), Inches(0.85),
            [{"runs": [("q = clamp(round(r/S) + Z, −128, 127)", 14.5, WHITE, True)],
              "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    textbox(slide, Inches(1.0), Inches(4.8), Inches(5.3), Inches(1.9),
            [{"runs": [("S", 12.5, CYAN, True), (" = hệ số tỷ lệ (Scale),  ", 12.5, WHITE, False),
                       ("Z", 12.5, CYAN, True), (" = điểm không (Zero-point), xác định trên mỗi layer từ tập hiệu chỉnh.", 12.5, WHITE, False)],
              "line": 1.1, "space_after": Pt(6)},
             {"runs": [("Giải lượng tử hóa khi suy luận:  r ≈ S·(q − Z)", 12.5, MUTED, False, True)]}])
    # cột phải: 3 lợi ích
    benefits = [
        ("Giảm bộ nhớ 4×", "File .tflite từ vài MB → ≈ 40 KB nhúng thẳng vào Flash.", GREEN),
        ("Tăng tốc tính toán", "Nhân ma trận số nguyên nhanh hơn → suy luận chỉ 216 ms.", CYAN),
        ("Tiết kiệm năng lượng", "Phép tính INT tiêu tốn ít điện hơn Float32.", AMBER),
    ]
    for i, (t, d, c) in enumerate(benefits):
        yy = Inches(2.5) + i * Inches(1.46)
        card(slide, Inches(6.75), yy, Inches(5.85), Inches(1.3), fill=CARD, line=c)
        number_badge(slide, Inches(7.0), yy + Inches(0.34), Inches(0.6), i + 1, c)
        textbox(slide, Inches(7.85), yy + Inches(0.18), Inches(4.6), Inches(0.45),
                [{"runs": [(t, 15, WHITE, True, False, FONT_B)]}])
        textbox(slide, Inches(7.85), yy + Inches(0.62), Inches(4.6), Inches(0.55),
                [{"runs": [(d, 12, MUTED, False)], "line": 1.0}])
    return slide


# ============================================================
# SLIDE 16 — EDGE AI Ý NGHĨA
# ============================================================
def slide_edge():
    slide = base_slide(3, "Cơ sở lý thuyết", "Ý nghĩa của Edge AI", accent=PURPLE,
                       subtitle="Khác biệt căn bản so với thiết bị IoT truyền thống chỉ thu thập rồi đẩy lên đám mây")
    feats = [
        ("Hoạt động khi mất mạng", "Vẫn phát hiện dị thường, chạy AI và kích hoạt còi/đèn dù WiFi mất hoàn toàn.", GREEN, "📴"),
        ("Độ trễ cực thấp", "Suy luận hoàn tất trong 216 ms ngay tại thiết bị, không có độ trễ mạng 100–500 ms.", CYAN, "⚡"),
        ("Bảo mật dữ liệu", "Dữ liệu xử lý cục bộ, chỉ gửi kết quả tổng hợp lên đám mây.", BLUE, "🔒"),
        ("Giảm chi phí Cloud", "Chỉ gửi 1 gói nhỏ mỗi 5 giây thay vì luồng dữ liệu thô liên tục.", AMBER, "💰"),
    ]
    cw = Inches(5.85)
    ch = Inches(1.9)
    x0 = Inches(0.72)
    y0 = Inches(2.95)
    gx = Inches(0.2)
    gy = Inches(0.2)
    for i, (t, d, c, ic) in enumerate(feats):
        x = x0 + (i % 2) * (cw + gx)
        y = y0 + (i // 2) * (ch + gy)
        card(slide, x, y, cw, ch, fill=CARD, line=c)
        circ = add_rect(slide, x + Inches(0.3), y + Inches(0.5), Inches(0.9), Inches(0.9),
                        CARD2, shape=MSO_SHAPE.OVAL, line=c, line_w=Pt(1.5))
        textbox(slide, x + Inches(0.3), y + Inches(0.5), Inches(0.9), Inches(0.9),
                [{"runs": [(ic, 24, c, False)], "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, x + Inches(1.4), y + Inches(0.35), cw - Inches(1.6), Inches(0.5),
                [{"runs": [(t, 16, WHITE, True, False, FONT_B)]}])
        textbox(slide, x + Inches(1.4), y + Inches(0.85), cw - Inches(1.6), Inches(0.9),
                [{"runs": [(d, 12.5, MUTED, False)], "line": 1.08}])
    return slide


# ============================================================
# SLIDE 17 — FIRMWARE LOOP
# ============================================================
def slide_firmware():
    slide = base_slide(4, "Triển khai Firmware", "Luồng hoạt động chính — loop() (chu kỳ 5 s)", accent=AMBER)
    steps = [
        ("Đọc cảm biến", "sensorManager.readAll() — BME280, Sharp, MQ135", BLUE),
        ("Lọc dị thường", "isAnomaly() — thay xₜ bằng EMAₜ₋₁ nếu Z > 3", PURPLE),
        ("Dịch cửa sổ + chuẩn hóa", "Shift buffer 24 bước, scaleData() → INT8", CYAN),
        ("Suy luận AI", "interpreter->Invoke() — LSTM INT8 (216 ms)", GREEN),
        ("Giải lượng tử + Inverse Scale", "Dequantize → final_pm25 (µg/m³)", TEAL),
        ("Cập nhật LCD", "displayManager.updateData() — Sprite DMA", BLUE),
        ("Kiểm tra cảnh báo", "final_pm25 ≥ 80 → Buzzer + LED ON", RED),
        ("Đồng bộ Firebase", "Realtime 5 s  •  History 60 s", AMBER),
    ]
    cols = 2
    cw = Inches(5.85)
    ch = Inches(0.92)
    x0 = Inches(0.72)
    y0 = Inches(2.5)
    gx = Inches(0.2)
    gy = Inches(0.16)
    for i, (t, d, c) in enumerate(steps):
        x = x0 + (i % cols) * (cw + gx)
        y = y0 + (i // cols) * (ch + gy)
        card(slide, x, y, cw, ch, fill=CARD)
        number_badge(slide, x + Inches(0.22), y + Inches(0.17), Inches(0.58), i + 1, c)
        textbox(slide, x + Inches(1.0), y + Inches(0.12), cw - Inches(1.2), Inches(0.42),
                [{"runs": [(t, 13.5, WHITE, True, False, FONT_B)]}])
        textbox(slide, x + Inches(1.0), y + Inches(0.5), cw - Inches(1.2), Inches(0.38),
                [{"runs": [(d, 11, MUTED, False)], "line": 1.0}])
    return slide


# ============================================================
# SLIDE 18 — WEB DASHBOARD
# ============================================================
def slide_web():
    slide = base_slide(5, "Web Dashboard", "Giao diện giám sát trực tuyến", accent=TEAL,
                       subtitle="aiot-air-quality.web.app  —  HTML5 / CSS3 / Vanilla JS + Chart.js v4, không framework nặng")
    # cột trái: đặc điểm
    card(slide, Inches(0.72), Inches(2.85), Inches(6.0), Inches(3.95), fill=CARD)
    textbox(slide, Inches(1.0), Inches(3.05), Inches(5.5), Inches(0.45),
            [{"runs": [("✨  Thiết kế Glassmorphism", 15, TEAL, True)]}])
    bullets(slide, Inches(1.0), Inches(3.6), Inches(5.5), Inches(3.1), [
        ("▸", [("Kính mờ ", 12.5, WHITE, False), ("backdrop-filter: blur(20px)", 12.5, CYAN, True),
               (" trên nền gradient tối.", 12.5, WHITE, False)]),
        ("▸", [("Màu theo ngưỡng PM2.5: ", 12.5, WHITE, False), ("xanh", 12.5, GREEN, True),
               (" / ", 12.5, WHITE, False), ("vàng", 12.5, AMBER, True), (" / ", 12.5, WHITE, False),
               ("cam", 12.5, RGBColor(0xFB,0x92,0x3C), True), (" / ", 12.5, WHITE, False),
               ("đỏ", 12.5, RED, True), (".", 12.5, WHITE, False)]),
        ("▸", [("Badge nhấp nháy ", 12.5, WHITE, False), ("[!] ANOMALY DETECTED", 12.5, RED, True),
               (" khi phát hiện gai nhiễu.", 12.5, WHITE, False)]),
        ("▸", [("Responsive Flexbox/Grid, breakpoint 768 px (Grid 2×2 trên mobile).", 12.5, WHITE, False)]),
    ], marker=TEAL, size=12.5, gap=Pt(10))
    # cột phải: 3 luồng dữ liệu + preload
    card(slide, Inches(6.92), Inches(2.85), Inches(5.7), Inches(3.95), fill=CARD)
    textbox(slide, Inches(7.22), Inches(3.05), Inches(5.2), Inches(0.45),
            [{"runs": [("📊  Ba luồng dữ liệu & tải trước", 15, BLUE, True)]}])
    lines = [
        ("PM2.5 thô", "raw — đọc trực tiếp từ cảm biến", RED),
        ("PM2.5 đã lọc", "filtered — sau Z-Score/EMA", BLUE),
        ("PM2.5 dự báo AI", "predicted — 1 giờ tới (LSTM)", GREEN),
    ]
    for i, (t, d, c) in enumerate(lines):
        yy = Inches(3.6) + i * Inches(0.72)
        add_rect(slide, Inches(7.22), yy + Inches(0.08), Inches(0.32), Inches(0.32), c,
                 shape=MSO_SHAPE.OVAL)
        textbox(slide, Inches(7.7), yy, Inches(4.7), Inches(0.7),
                [{"runs": [(t + "  ", 13, WHITE, True), (d, 11.5, MUTED, False, True)]}],
                anchor=MSO_ANCHOR.MIDDLE)
    band = add_rect(slide, Inches(7.22), Inches(5.85), Inches(5.1), Inches(0.78),
                    CARD2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=TEAL, radius=0.1)
    textbox(slide, Inches(7.4), Inches(5.85), Inches(4.8), Inches(0.78),
            [{"runs": [("Tải trước 30 bản ghi lịch sử ", 12, WHITE, False),
                       ("(limitToLast(30))", 12, CYAN, True),
                       (" → biểu đồ liền mạch ngay khi mở trang.", 12, WHITE, False)], "line": 1.05}],
            anchor=MSO_ANCHOR.MIDDLE)
    return slide


# ============================================================
# SLIDE 19 — KẾT QUẢ HIỆU NĂNG
# ============================================================
def slide_perf():
    slide = base_slide(6, "Kết quả & Kết luận", "Kết quả đo đạc hiệu năng", accent=RED)
    # KPI lớn
    kpis = [
        ("216 ms", "Suy luận LSTM INT8", CYAN),
        ("4,3 %", "CPU dành cho AI", GREEN),
        ("≈ 40 KB", "Mô hình INT8 (Flash)", AMBER),
        ("506,6", "µg/m³ gai nhiễu bị chặn", RED),
    ]
    cw = Inches(2.92)
    ch = Inches(1.5)
    x0 = Inches(0.72)
    y0 = Inches(2.45)
    gx = Inches(0.18)
    for i, (v, d, c) in enumerate(kpis):
        x = x0 + i * (cw + gx)
        card(slide, x, y0, cw, ch, fill=CARD, line=c)
        textbox(slide, x, y0 + Inches(0.2), cw, Inches(0.7),
                [{"runs": [(v, 28, c, True, False, FONT_B)], "align": PP_ALIGN.CENTER}])
        textbox(slide, x, y0 + Inches(0.95), cw, Inches(0.45),
                [{"runs": [(d, 11.5, MUTED, False)], "align": PP_ALIGN.CENTER}])
    # bảng chi tiết
    rows = [
        ("Chỉ số đo lường", "Giá trị", "Ghi chú"),
        ("Thời gian suy luận LSTM INT8", "216 ms", "Đo bằng millis()"),
        ("Chu kỳ vòng lặp chính", "5000 ms", "delay(5000)"),
        ("Tỉ lệ CPU dành cho AI", "4,3 %", "216 / 5000"),
        ("Mô hình TFLite (INT8) / gốc (Float32)", "≈ 40 / 160 KB", "Giảm 4× sau lượng tử hóa"),
        ("Tensor Arena trong PSRAM", "2048 KB", "heap_caps_malloc"),
        ("Chu kỳ realtime / history", "5 s / 60 s", "/sensor,/ai  •  /history"),
    ]
    add_table(slide, rows, Inches(0.72), Inches(4.2), Inches(11.9), Inches(2.55),
              col_w=[Inches(5.6), Inches(3.0), Inches(3.3)], header_color=RED,
              font_size=12, header_size=12.5)
    return slide


# ============================================================
# SLIDE 20 — THỬ NGHIỆM BỘ LỌC (CHART)
# ============================================================
def slide_experiment():
    slide = base_slide(6, "Kết quả & Kết luận", "Thực nghiệm bộ lọc dị thường Z-Score", accent=RED)
    # hình chart bên trái
    card(slide, Inches(0.72), Inches(2.5), Inches(7.3), Inches(4.3), fill=CARD)
    chart = os.path.join(MEDIA, "chart.png")
    if os.path.exists(chart):
        fit_image(slide, chart, Inches(0.92), Inches(2.7), Inches(6.9), Inches(3.9))
    # quan sát bên phải
    card(slide, Inches(8.2), Inches(2.5), Inches(4.42), Inches(4.3), fill=CARD, line=RED)
    textbox(slide, Inches(8.5), Inches(2.7), Inches(3.9), Inches(0.5),
            [{"runs": [("Kịch bản: thổi khói trực tiếp", 14, RED, True)]}])
    bullets(slide, Inches(8.5), Inches(3.3), Inches(3.95), Inches(3.4), [
        ("▸", [("PM2.5 thô vọt lên ", 12.5, WHITE, False), ("506,6 µg/m³", 12.5, RED, True),
               (".", 12.5, WHITE, False)]),
        ("▸", [("Serial báo: ", 12.5, WHITE, False),
               ("[ANOMALY] Filtered out!", 11.5, AMBER, True)]),
        ("▸", [("Đường PM2.5 đã lọc vẫn ", 12.5, WHITE, False), ("mượt mà", 12.5, GREEN, True),
               (" — gai nhiễu bị chặn 100%.", 12.5, WHITE, False)]),
        ("▸", [("Dự báo AI ", 12.5, WHITE, False), ("không bị ảnh hưởng", 12.5, GREEN, True),
               (" vì nhận giá trị EMA ổn định.", 12.5, WHITE, False)]),
    ], marker=RED, size=12.5, gap=Pt(11))
    return slide


# ============================================================
# SLIDE 21 — SẢN PHẨM + CẢNH BÁO SỚM
# ============================================================
def slide_product():
    slide = base_slide(6, "Kết quả & Kết luận", "Sản phẩm thực tế & cảnh báo sớm", accent=RED)
    # ảnh sản phẩm (xoay 90 độ)
    card(slide, Inches(0.72), Inches(2.5), Inches(6.0), Inches(4.3), fill=CARD)
    prod = os.path.join(MEDIA, "product.jpg")
    rotated = os.path.join(MEDIA, "_product_rot.png")
    try:
        im = Image.open(prod).rotate(-90, expand=True)
        im.save(rotated)
        prod_use = rotated
    except Exception:
        prod_use = prod
    if os.path.exists(prod_use):
        fit_image(slide, prod_use, Inches(0.92), Inches(2.7), Inches(5.6), Inches(3.9))
    # cột phải: logic cảnh báo sớm
    card(slide, Inches(6.92), Inches(2.5), Inches(5.7), Inches(4.3), fill=CARD, line=AMBER)
    textbox(slide, Inches(7.22), Inches(2.7), Inches(5.2), Inches(0.5),
            [{"runs": [("🔔  Cảnh báo dựa trên DỰ BÁO AI", 15, AMBER, True)]}])
    textbox(slide, Inches(7.22), Inches(3.25), Inches(5.15), Inches(0.95),
            [{"runs": [("Hệ thống phản ứng ", 13, WHITE, False), ("trước", 13, AMBER, True),
                       (" khi ô nhiễm đạt mức nguy hiểm — cho người dùng thời gian đeo khẩu trang, bật máy lọc, đóng cửa sổ.", 13, WHITE, False)], "line": 1.12}])
    eq = add_rect(slide, Inches(7.22), Inches(4.35), Inches(5.1), Inches(0.8),
                  CARD2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=RED, radius=0.1)
    textbox(slide, Inches(7.22), Inches(4.35), Inches(5.1), Inches(0.8),
            [{"runs": [("final_pm25 ≥ 80 µg/m³  →  Buzzer + LED", 13.5, WHITE, True)],
              "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    bullets(slide, Inches(7.22), Inches(5.4), Inches(5.15), Inches(1.3), [
        ("▸", [("Ngưỡng 80 µg/m³ = mức ", 12, WHITE, False),
               ("\"Unhealthy for Sensitive Groups\"", 12, AMBER, True), (" (EPA AQI).", 12, WHITE, False)]),
        ("▸", [("Còi 1000 Hz không chặn luồng (non-blocking).", 12, WHITE, False)]),
    ], marker=AMBER, size=12, gap=Pt(7))
    return slide


# ============================================================
# SLIDE 22 — ĐÁNH GIÁ & HẠN CHẾ
# ============================================================
def slide_eval():
    slide = base_slide(6, "Kết quả & Kết luận", "Đánh giá tổng thể & điểm hạn chế", accent=RED)
    card(slide, Inches(0.72), Inches(2.5), Inches(6.0), Inches(4.3), fill=CARD, line=GREEN)
    textbox(slide, Inches(1.02), Inches(2.7), Inches(5.5), Inches(0.5),
            [{"runs": [("✔  Đã đạt được", 16, GREEN, True)]}])
    bullets(slide, Inches(1.02), Inches(3.3), Inches(5.45), Inches(3.3), [
        ("●", [("Hoạt động ổn định 24/7, độ tin cậy cao.", 13, WHITE, False)]),
        ("●", [("Bộ lọc Z-Score+EMA chặn 100% gai nhiễu trong kiểm thử.", 13, WHITE, False)]),
        ("●", [("LSTM INT8 chạy hiệu quả, suy luận 216 ms.", 13, WHITE, False)]),
        ("●", [("Web Dashboard responsive, giám sát từ xa hữu ích.", 13, WHITE, False)]),
    ], marker=GREEN, size=13, gap=Pt(11))
    card(slide, Inches(6.92), Inches(2.5), Inches(5.7), Inches(4.3), fill=CARD, line=AMBER)
    textbox(slide, Inches(7.22), Inches(2.7), Inches(5.2), Inches(0.5),
            [{"runs": [("⚠  Hạn chế hiện tại", 16, AMBER, True)]}])
    bullets(slide, Inches(7.22), Inches(3.3), Inches(5.15), Inches(3.3), [
        ("●", [("Sharp GP2Y1014 cho chỉ số ", 13, WHITE, False), ("tương đối", 13, AMBER, True),
               (", chưa hiệu chuẩn theo PMS7003.", 13, WHITE, False)]),
        ("●", [("LSTM huấn luyện trên ", 13, WHITE, False), ("dữ liệu giả lập", 13, AMBER, True),
               (", chưa đủ dữ liệu thực nhiều tuần.", 13, WHITE, False)]),
        ("●", [("Chưa tích hợp pin sạc — phụ thuộc nguồn USB cố định.", 13, WHITE, False)]),
    ], marker=AMBER, size=13, gap=Pt(11))
    return slide


# ============================================================
# SLIDE 23 — KẾT LUẬN & HƯỚNG PHÁT TRIỂN
# ============================================================
def slide_future():
    slide = base_slide(6, "Kết quả & Kết luận", "Kết luận & hướng phát triển", accent=RED)
    textbox(slide, Inches(0.72), Inches(2.25), Inches(11.9), Inches(0.5),
            [{"runs": [("Hoàn thành cả 4 mục tiêu: phần cứng đa chỉ số • lọc nhiễu thông minh • Edge AI tại biên • hệ sinh thái Cloud.",
                        13.5, TEAL, False, True)]}])
    futures = [
        ("Nâng cấp cảm biến PMS7003", "Cảm biến laser UART, PM1.0/2.5/10 hiệu chuẩn ±10 µg/m³.", CYAN),
        ("Thu thập dữ liệu thực & tái huấn luyện", "4–6 tuần dữ liệu Việt Nam, thêm tầng Attention cho LSTM.", PURPLE),
        ("Cập nhật firmware OTA", "ArduinoOTA / Firebase Storage cập nhật từ xa qua WiFi.", BLUE),
        ("Thiết kế PCB & tích hợp pin", "Bo mạch in gọn, sạc Li-Po, Deep Sleep tiết kiệm điện.", GREEN),
        ("Mở rộng mạng lưới trạm đo", "Nhiều thiết bị + heatmap phân bố ô nhiễm theo không gian.", AMBER),
    ]
    x0 = Inches(0.72)
    y0 = Inches(2.95)
    cw = Inches(11.9)
    ch = Inches(0.72)
    gy = Inches(0.11)
    for i, (t, d, c) in enumerate(futures):
        y = y0 + i * (ch + gy)
        card(slide, x0, y, cw, ch, fill=CARD)
        number_badge(slide, x0 + Inches(0.24), y + Inches(0.08), Inches(0.55), i + 1, c)
        textbox(slide, x0 + Inches(1.05), y, Inches(4.6), ch,
                [{"runs": [(t, 14, WHITE, True, False, FONT_B)]}], anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, x0 + Inches(5.7), y, Inches(6.0), ch,
                [{"runs": [(d, 12, MUTED, False)]}], anchor=MSO_ANCHOR.MIDDLE)
    return slide


# ============================================================
# SLIDE 24 — CẢM ƠN
# ============================================================
def slide_thanks():
    slide = prs.slides.add_slide(BLANK)
    gradient_bg(slide)
    add_rect(slide, SW - Inches(4.3), Inches(-1.8), Inches(6), Inches(6),
             CARD2, shape=MSO_SHAPE.OVAL)
    add_rect(slide, Inches(-1.6), SH - Inches(3.2), Inches(5), Inches(5),
             CARD, shape=MSO_SHAPE.OVAL)
    add_rect(slide, 0, 0, Inches(0.18), SH, CYAN)
    textbox(slide, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.4),
            [{"runs": [("Cảm ơn Thầy và các bạn", 44, WHITE, True, False, FONT_B)], "space_after": Pt(4)},
             {"runs": [("đã lắng nghe!", 44, CYAN, True, False, FONT_B)]}])
    add_rect(slide, Inches(0.95), Inches(4.35), Inches(1.6), Pt(3.5), CYAN,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    textbox(slide, Inches(0.95), Inches(4.65), Inches(11), Inches(0.6),
            [{"runs": [("Q & A  —  Rất mong nhận được góp ý từ Thầy và các bạn", 16, MUTED, False, True)]}])
    # thông tin liên hệ / demo
    card(slide, Inches(0.95), Inches(5.6), Inches(11.4), Inches(1.0), fill=CARD, line=LINE)
    textbox(slide, Inches(1.25), Inches(5.6), Inches(11.0), Inches(1.0),
            [{"runs": [("🌐  Demo trực tuyến:  ", 14, TEAL, True),
                       ("https://aiot-air-quality.web.app", 14, WHITE, False),
                       ("        •        Nhóm KANT  —  23CT111", 13, MUTED, False)]}],
            anchor=MSO_ANCHOR.MIDDLE)
    return slide


# ---------------- BUILD ----------------
slide_cover()
slide_agenda()
slide_problem()
slide_solution()
slide_objectives()
slide_scope()
slide_architecture()
slide_bom()
slide_pinout()
slide_pipeline()
slide_scaler()
slide_zscore()
slide_lstm()
slide_model()
slide_quant()
slide_edge()
slide_firmware()
slide_web()
slide_perf()
slide_experiment()
slide_product()
slide_eval()
slide_future()
slide_thanks()

prs.save(OUT)
print("Saved:", OUT, "| slides:", len(prs.slides._sldIdLst))
